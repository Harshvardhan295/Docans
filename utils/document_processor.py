# utils/document_processor.py
import io
import fitz  # PyMuPDF
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter


# 1000 chunk size and 200 overlap is an industry-standard starting point for RAG
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200,
    length_function=len,
    is_separator_regex=False,
)

def extract_text_from_pdf(file_bytes: bytes) -> str:
    text = ""
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page_num, page in enumerate(doc):
            page_text = page.get_text("text")
            if page_text.strip():
                text += f"\n--- Page {page_num + 1} ---\n{page_text}"
    return text

def extract_text_from_pptx(file_bytes: bytes) -> str:
    text = ""
    presentation = Presentation(io.BytesIO(file_bytes))
    for slide_num, slide in enumerate(presentation.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        if slide_text.strip():
            text += f"\n--- Slide {slide_num + 1} ---\n{slide_text}"
    return text

def chunk_document(text: str) -> list[str]:
    if not text.strip():
        return []
    chunks = text_splitter.split_text(text)
    return chunks

async def process_file(file_name: str, file_bytes: bytes) -> dict:
    extracted_text = ""
    
    if file_name.lower().endswith('.pdf'):
        extracted_text = extract_text_from_pdf(file_bytes)
    elif file_name.lower().endswith('.pptx'):
        extracted_text = extract_text_from_pptx(file_bytes)
    else:
        raise ValueError("Unsupported file format. Please upload PDF or PPTX.")

    chunks = chunk_document(extracted_text)
    
    return {
        "filename": file_name,
        "total_characters": len(extracted_text),
        "total_chunks": len(chunks),
        "chunks": chunks
    }