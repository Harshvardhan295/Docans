import io
import re
import fitz  # PyMuPDF
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=500,
    chunk_overlap=100,
    length_function=len,
)


def clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def extract_text_from_pdf(file_bytes: bytes) -> list[dict]:
    sections = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for index, page in enumerate(doc, start=1):
            text = clean_text(page.get_text("text"))
            if text:
                sections.append({"source": str(index), "text": text})
    return sections


def extract_text_from_pptx(file_bytes: bytes) -> list[dict]:
    sections = []
    presentation = Presentation(io.BytesIO(file_bytes))
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text = clean_text("\n".join(slide_text))
        if text:
            sections.append({"source": str(index), "text": text})
    return sections


async def process_file(file_name: str, file_bytes: bytes) -> dict:
    if file_name.lower().endswith(".pdf"):
        sections = extract_text_from_pdf(file_bytes)
    elif file_name.lower().endswith(".pptx"):
        sections = extract_text_from_pptx(file_bytes)
    else:
        raise ValueError("Unsupported format.")

    chunks = []
    total_characters = 0

    for section in sections:
        total_characters += len(section["text"])
        for chunk in text_splitter.split_text(section["text"]):
            chunks.append({"text": chunk, "source": section["source"]})

    return {
        "filename": file_name,
        "total_characters": total_characters,
        "total_chunks": len(chunks),
        "chunks": chunks,
    }
